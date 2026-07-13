"""Document parsers for multiple file types"""

import os
import re
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def parse_document(file_path: str) -> str:
    """Parse a document and return its plain text content.
    Supports PDF, DOCX, PPTX, XLSX, MD/TXT, HTML.
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    parsers = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".doc": _parse_docx,
        ".pptx": _parse_pptx,
        ".ppt": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".xls": _parse_xlsx,
        ".md": _parse_text,
        ".txt": _parse_text,
        ".html": _parse_html,
        ".htm": _parse_html,
    }

    parser = parsers.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {ext}")

    logger.info(f"Parsing: {file_path} ({ext})")
    return parser(file_path)


def _parse_pdf(path: str) -> str:
    import fitz
    doc = fitz.open(path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            pages.append(f"[Page {i+1}]\n{text}")
    doc.close()
    return "\n\n".join(pages)


def _parse_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras)


def _parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        if texts:
            slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _parse_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            line = " | ".join(str(c) for c in row if c is not None)
            if line.strip():
                rows.append(line)
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


def _parse_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _parse_html(path: str) -> str:
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)
